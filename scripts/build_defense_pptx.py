"""
박사학위 논문 방어 발표 PPTX 자동 생성기 — 초안 (draft).

Shin Dong-wook · Hoseo University · 2026
논문: "근위 정책 최적화 기반 적응형 동적 가중치 학습을 통한 Triple-Hybrid RAG
프레임워크의 성능 최적화 연구"

Usage:
    python scripts/build_defense_pptx.py \
        --output thesis_current/박사논문_방어발표_v1.pptx

Design goals:
- 16:9 와이드스크린, ~26 slides (15~20분 발표)
- Part A (Opening) → B (문제·방법) → C (결과) → D (결론) → Appendix
- 논문 수치를 하드코딩 (후속 수정 시 이 파일에서만 관리)
- docs/figures/*.png 9개 그림 활용
- 한국어 폰트: 맑은 고딕 (한글), Arial (수식/영문)
"""

from __future__ import annotations

import argparse
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

ROOT = Path(__file__).resolve().parent.parent
FIGS = ROOT / "docs" / "figures"

# Hoseo 대학 블루 + 강조 색상
HOSEO_BLUE = RGBColor(0x1E, 0x3A, 0x8A)
HOSEO_BLUE_LIGHT = RGBColor(0x4F, 0x86, 0xC6)
ACCENT_RED = RGBColor(0xE8, 0x75, 0x6E)
ACCENT_GREEN = RGBColor(0x8F, 0xB5, 0x73)
TEXT_DARK = RGBColor(0x1E, 0x29, 0x3B)
TEXT_MUTED = RGBColor(0x6B, 0x72, 0x80)
BG_SOFT = RGBColor(0xF5, 0xF7, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

KO_FONT = "맑은 고딕"
EN_FONT = "Arial"

SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)


# --------------------------------------------------------------------------
# primitives
# --------------------------------------------------------------------------

def _set_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, left, top, width, height, text: str,
              size: int = 18, bold: bool = False,
              color: RGBColor = TEXT_DARK, align=PP_ALIGN.LEFT,
              font: str = KO_FONT, anchor=MSO_ANCHOR.TOP) -> None:
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color


def _add_rect(slide, left, top, width, height,
              fill_color: RGBColor = HOSEO_BLUE,
              line_color: RGBColor | None = None) -> None:
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
    return shp


def _add_rounded_card(slide, left, top, width, height,
                      title: str, body: str,
                      title_color: RGBColor = HOSEO_BLUE,
                      body_color: RGBColor = TEXT_DARK,
                      fill_color: RGBColor = WHITE,
                      border: RGBColor = HOSEO_BLUE_LIGHT) -> None:
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.color.rgb = border
    shp.line.width = Pt(1.2)
    shp.shadow.inherit = False
    tf = shp.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
    # Title
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.LEFT
    r0 = p0.add_run()
    r0.text = title
    r0.font.name = KO_FONT
    r0.font.size = Pt(16)
    r0.font.bold = True
    r0.font.color.rgb = title_color
    # Body
    for line in body.split("\n"):
        if not line.strip():
            continue
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = line
        r.font.name = KO_FONT
        r.font.size = Pt(12)
        r.font.color.rgb = body_color


def _add_footer(slide, slide_no: int, total: int,
                section: str = "") -> None:
    # Thin blue bar at bottom
    _add_rect(slide, 0, SLIDE_H - Inches(0.35), SLIDE_W, Inches(0.05),
              fill_color=HOSEO_BLUE)
    _add_text(
        slide,
        Inches(0.5), SLIDE_H - Inches(0.3),
        SLIDE_W - Inches(1), Inches(0.3),
        f"신동욱 (Shin Dong-wook) · 호서대학교 대학원 융합공학과 · 박사학위 방어 발표",
        size=10, color=TEXT_MUTED, align=PP_ALIGN.LEFT,
    )
    _add_text(
        slide,
        SLIDE_W - Inches(1.2), SLIDE_H - Inches(0.3),
        Inches(0.8), Inches(0.3),
        f"{slide_no} / {total}",
        size=10, color=TEXT_MUTED, align=PP_ALIGN.RIGHT,
    )
    if section:
        _add_text(
            slide,
            SLIDE_W - Inches(3.5), SLIDE_H - Inches(0.3),
            Inches(2.0), Inches(0.3),
            section,
            size=10, color=HOSEO_BLUE, align=PP_ALIGN.RIGHT, bold=True,
        )


def _add_section_header(slide, section_kr: str, section_en: str) -> None:
    # Top-left accent bar + title line
    _add_rect(slide, 0, 0, Inches(0.25), SLIDE_H, fill_color=HOSEO_BLUE)
    _add_text(
        slide, Inches(0.5), Inches(0.3),
        SLIDE_W - Inches(1), Inches(0.6),
        section_kr,
        size=28, bold=True, color=HOSEO_BLUE,
    )
    _add_text(
        slide, Inches(0.5), Inches(0.85),
        SLIDE_W - Inches(1), Inches(0.4),
        section_en,
        size=13, color=TEXT_MUTED, font=EN_FONT,
    )
    # underline
    _add_rect(slide, Inches(0.5), Inches(1.35),
              Inches(1.2), Emu(25000), fill_color=HOSEO_BLUE_LIGHT)


# --------------------------------------------------------------------------
# slide builders
# --------------------------------------------------------------------------

def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def build_s01_title(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    # Blue left band
    _add_rect(slide, 0, 0, Inches(4.5), SLIDE_H, fill_color=HOSEO_BLUE)
    # white on blue: school + label
    _add_text(
        slide, Inches(0.5), Inches(0.7),
        Inches(3.8), Inches(0.5),
        "호서대학교 대학원",
        size=20, bold=True, color=WHITE,
    )
    _add_text(
        slide, Inches(0.5), Inches(1.1),
        Inches(3.8), Inches(0.4),
        "융합공학과",
        size=14, color=WHITE,
    )
    _add_text(
        slide, Inches(0.5), Inches(5.7),
        Inches(3.8), Inches(0.4),
        "박사학위 논문 방어 발표",
        size=14, bold=True, color=WHITE,
    )
    _add_text(
        slide, Inches(0.5), Inches(6.1),
        Inches(3.8), Inches(0.4),
        "2026. 4. 30.",
        size=12, color=WHITE,
    )

    # Right side: title + author
    _add_text(
        slide, Inches(5.0), Inches(1.4),
        Inches(7.8), Inches(0.5),
        "박사학위 논문",
        size=14, color=TEXT_MUTED,
    )
    _add_text(
        slide, Inches(5.0), Inches(1.9),
        Inches(8.0), Inches(2.4),
        "근위 정책 최적화 기반 적응형\n동적 가중치 학습을 통한\nTriple-Hybrid RAG 프레임워크의\n성능 최적화 연구",
        size=28, bold=True, color=HOSEO_BLUE,
    )
    _add_text(
        slide, Inches(5.0), Inches(4.5),
        Inches(8.0), Inches(0.6),
        "Performance Optimization of Triple-Hybrid RAG Framework\nvia Proximal Policy Optimization-based Learned Dynamic Weighting",
        size=12, color=TEXT_MUTED, font=EN_FONT,
    )
    # author block
    _add_text(
        slide, Inches(5.0), Inches(5.7),
        Inches(8.0), Inches(0.4),
        "발표자: 신동욱 (Shin Dong-wook)",
        size=16, bold=True, color=TEXT_DARK,
    )
    _add_text(
        slide, Inches(5.0), Inches(6.1),
        Inches(8.0), Inches(0.4),
        "지도교수: 문남미 (Prof. Nammee Moon)",
        size=14, color=TEXT_DARK,
    )


def build_s02_agenda(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "발표 목차", "Agenda")

    sections = [
        ("Part A", "Opening", "연구 한눈에 보기 · 핵심 결과 요약", HOSEO_BLUE),
        ("Part B", "문제와 방법", "Triple-Hybrid RAG · α·β·γ 가중치 문제 · R-DWA 와 L-DWA",
         HOSEO_BLUE_LIGHT),
        ("Part C", "결과", "5,000 QA 정량 비교 · Oracle 초과 · Conditional 개선",
         ACCENT_RED),
        ("Part D", "결론", "기여 · 한계 · 향후 연구 방향", ACCENT_GREEN),
    ]

    top = Inches(2.0)
    for tag, kr, desc, color in sections:
        _add_rect(slide, Inches(0.8), top, Inches(0.2), Inches(0.8),
                  fill_color=color)
        _add_text(
            slide, Inches(1.1), top,
            Inches(2.0), Inches(0.4),
            tag,
            size=12, bold=True, color=color,
        )
        _add_text(
            slide, Inches(1.1), top + Inches(0.35),
            Inches(3.0), Inches(0.5),
            kr,
            size=20, bold=True, color=TEXT_DARK,
        )
        _add_text(
            slide, Inches(5.0), top + Inches(0.15),
            Inches(7.5), Inches(0.7),
            desc,
            size=14, color=TEXT_MUTED,
        )
        top += Inches(1.05)

    _add_text(
        slide, Inches(0.8), Inches(6.6),
        Inches(12), Inches(0.3),
        "발표 시간: 약 18분 · 슬라이드 26장 + Appendix",
        size=11, color=TEXT_MUTED, align=PP_ALIGN.RIGHT,
    )


def build_s03_summary(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "한 문장 요약", "Research in One Sentence")

    _add_text(
        slide, Inches(0.8), Inches(2.0),
        Inches(11.8), Inches(1.6),
        "Triple-Hybrid RAG 에서 Vector · Graph · Ontology 세 지식 소스의\n"
        "가중치 결정을 PPO 로 학습한 연속 정책 (L-DWA) 이 이산 격자 Oracle\n"
        "상한을 미세하게 초과하는 지점에 도달하였다.",
        size=22, bold=True, color=HOSEO_BLUE,
    )

    # Key 3 findings as cards
    top = Inches(4.2)
    w = Inches(3.9)
    gap = Inches(0.15)
    cards = [
        ("① Oracle 초과",
         "L-DWA 3-seed 평균 F1_strict 0.562\n"
         "> 이산 격자 Oracle 0.554 (+0.8%p)\n"
         "네 개 F1 지표 모두에서 확인\n"
         "연속 Dirichlet 평균이 격자 바깥 활용",
         HOSEO_BLUE),
        ("② Conditional +36.7%",
         "조건부 질의 (n=1,250) 에서\n"
         "R-DWA 0.223 → L-DWA 0.304\n"
         "Oracle (0.290) 마저 초과\n"
         "유형 내 미세 구분을 학습이 포착",
         ACCENT_RED),
        ("③ 학습 비용 −85%",
         "오프라인 보상 캐시 (330K, 16.8MB)\n"
         "로 학습 중 LLM 호출 0 회\n"
         "$288 → $33 · 3-seed 5시간\n"
         "개인 연구자 수준 재현 가능",
         ACCENT_GREEN),
    ]
    left = Inches(0.8)
    for title, body, color in cards:
        _add_rounded_card(
            slide, left, top, w, Inches(2.5),
            title=title, body=body, title_color=color,
        )
        left += w + gap


def build_s04_rag_bg(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "배경 — RAG 란 무엇인가", "Background · Retrieval-Augmented Generation")

    _add_text(
        slide, Inches(0.8), Inches(2.0),
        Inches(7.5), Inches(2.0),
        "대형 언어모델 (LLM) 은 내부 지식만으로 답하면\n"
        "① 환각 (없는 사실을 지어냄),\n"
        "② 최신성 부족,\n"
        "③ 출처 불명 문제가 발생한다.\n\n"
        "RAG 는 답을 만들기 전에 먼저 관련 문서를 검색하여\n"
        "LLM 에 참고 맥락으로 제공하는 기법이다.",
        size=15, color=TEXT_DARK,
    )

    # Simple pipeline diagram
    _add_rounded_card(slide, Inches(8.6), Inches(2.2), Inches(4.3), Inches(0.9),
                      title="① 질의 (Query)", body="사용자 질문",
                      fill_color=BG_SOFT)
    _add_rounded_card(slide, Inches(8.6), Inches(3.3), Inches(4.3), Inches(0.9),
                      title="② Retrieval (검색)", body="지식 소스에서 관련 정보 추출",
                      fill_color=BG_SOFT, title_color=ACCENT_RED)
    _add_rounded_card(slide, Inches(8.6), Inches(4.4), Inches(4.3), Inches(0.9),
                      title="③ Generation (응답)", body="LLM 이 맥락을 바탕으로 응답",
                      fill_color=BG_SOFT, title_color=ACCENT_GREEN)

    _add_text(
        slide, Inches(0.8), Inches(5.5),
        Inches(12), Inches(1.2),
        "본 논문의 관심은 ② Retrieval 단계. 특히 Vector · Graph · Ontology\n"
        "세 이질적 소스를 동시에 사용하는 Triple-Hybrid 구조에서 각 소스의\n"
        "가중치 (α, β, γ) 를 어떻게 정할 것인가가 연구 질문이다.",
        size=13, color=HOSEO_BLUE, bold=True,
    )


def build_s05_triple_hybrid(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "Triple-Hybrid — 세 지식 소스", "Three Heterogeneous Knowledge Sources")

    top = Inches(2.0)
    w = Inches(3.9)
    gap = Inches(0.15)
    sources = [
        ("🔵 Vector 검색",
         "문장 임베딩 유사도\n"
         "FAISS, top-k=3\n\n"
         "강점:\n• 의미 기반 매칭\n• 어휘 변동에 유연\n\n약점:\n• 정확한 관계 추론 어려움",
         HOSEO_BLUE),
        ("🟢 Graph 검색",
         "엔티티-관계 BFS\n"
         "NetworkX, depth=3\n\n"
         "강점:\n• 멀티홉 관계 질의\n• 소속·참여 정보\n\n약점:\n• 유사 표현에 취약",
         ACCENT_GREEN),
        ("🔴 Ontology 추론",
         "Owlready2 + HermiT\n"
         "클래스·속성 추론\n\n"
         "강점:\n• 조건부 제약 (나이·자격)\n• 논리적 만족성\n\n약점:\n• 스키마에 의존",
         ACCENT_RED),
    ]
    left = Inches(0.8)
    for title, body, color in sources:
        _add_rounded_card(
            slide, left, top, w, Inches(4.2),
            title=title, body=body, title_color=color,
        )
        left += w + gap

    _add_text(
        slide, Inches(0.8), Inches(6.4),
        Inches(11.5), Inches(0.7),
        "세 소스는 서로 강점이 다름 → 질의마다 적절한 비율로 섞어야 함.",
        size=15, color=HOSEO_BLUE, bold=True, align=PP_ALIGN.CENTER,
    )


def build_s06_weight_problem(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "α·β·γ — 가중치 문제", "The Weighting Problem")

    _add_text(
        slide, Inches(0.8), Inches(2.0),
        Inches(6), Inches(0.6),
        "세 소스의 기여도를 α·β·γ 로 표기, 합 = 1 제약",
        size=16, bold=True, color=TEXT_DARK,
    )
    _add_text(
        slide, Inches(0.8), Inches(2.6),
        Inches(6), Inches(2.5),
        "→ (α, β, γ) 는 Δ³ 삼각형 안의 한 점\n\n"
        "예시:\n"
        "• (0.6, 0.2, 0.2) — Vector 중심\n"
        "• (0.2, 0.6, 0.2) — Graph 중심\n"
        "• (0.2, 0.2, 0.6) — Ontology 중심\n"
        "• (1/3, 1/3, 1/3) — 균등",
        size=14, color=TEXT_DARK,
    )

    # Use figure 4-1 for visual
    if (FIGS / "fig4_1_rdwa_vs_oracle.png").exists():
        slide.shapes.add_picture(
            str(FIGS / "fig4_1_rdwa_vs_oracle.png"),
            Inches(7.2), Inches(1.9),
            width=Inches(5.8),
        )

    _add_text(
        slide, Inches(0.8), Inches(5.7),
        Inches(11.7), Inches(1.2),
        "핵심 어려움 — 질의마다 최적의 점이 다름.  \n"
        "• 「최재원 교수의 소속 학과는?」 → Vector 중심  \n"
        "• 「55세 이하 소속 교수는?」 → Ontology 중심  \n"
        "• 「배 포함 프로젝트 참여 학과는?」 → Graph 중심",
        size=13, color=TEXT_MUTED,
    )


def build_s07_rdwa(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "선행 연구 — R-DWA (규칙 기반)", "Prior Work · Rule-based Dynamic Weighting")

    _add_text(
        slide, Inches(0.8), Inches(2.0),
        Inches(6), Inches(0.5),
        "Shin & Moon (2025, JKSCI) — 2단계 규칙",
        size=16, bold=True, color=HOSEO_BLUE,
    )

    # Table mock
    _add_rounded_card(
        slide, Inches(0.8), Inches(2.6), Inches(6), Inches(2.3),
        title="Stage 1-2. 유형별 기본 가중치 (Table 4-1)",
        body="유형             α       β       γ\n"
             "simple           0.6     0.2     0.2\n"
             "multi_hop        0.2     0.6     0.2\n"
             "conditional      0.2     0.2     0.6\n\n"
             "+ 밀도 신호 s_r, s_c 로 미세 조정 (λ = 0.3)",
    )

    # 4 limits
    _add_text(
        slide, Inches(7.2), Inches(2.0),
        Inches(5.8), Inches(0.5),
        "R-DWA 의 네 가지 구조적 한계",
        size=16, bold=True, color=ACCENT_RED,
    )
    _add_text(
        slide, Inches(7.2), Inches(2.6),
        Inches(5.8), Inches(4.0),
        "① 도메인 편향된 기본 가중치\n"
        "    — Oracle 실측은 Ontology 지배적\n\n"
        "② 단일 λ = 0.3 하이퍼파라미터\n"
        "    — 질의별 특성 미반영\n\n"
        "③ 유형 내 미세 구분 부재\n"
        "    — 같은 conditional 도 성격 다양\n\n"
        "④ 도메인 이전 불가\n"
        "    — HotpotQA F1 0.097 < Vector-only 0.102",
        size=13, color=TEXT_DARK,
    )

    _add_text(
        slide, Inches(0.8), Inches(6.5),
        Inches(11.7), Inches(0.7),
        "→ 본 논문은 이 네 한계를 정면으로 다루기 위해 규칙 표를 학습된 신경망으로 대체한다.",
        size=13, color=HOSEO_BLUE, bold=True, align=PP_ALIGN.CENTER,
    )


def build_s08_ldwa_mdp(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "제안 — L-DWA (MDP 공식화)", "Proposed · Learned DWA as 1-step MDP")

    _add_text(
        slide, Inches(0.8), Inches(2.0),
        Inches(11.7), Inches(0.5),
        "가중치 결정을 Markov Decision Process 로 환원 — RAG 맥락에서 최초 사례",
        size=16, bold=True, color=HOSEO_BLUE,
    )

    cards = [
        ("🧩 State (18-dim)",
         "• 질의 메타 (길이·엔티티·제약)\n"
         "• Intent logits (3-dim)\n"
         "• 밀도 신호 (s_e, s_r, s_c)\n"
         "• Retrieval 점수 통계량",
         HOSEO_BLUE),
        ("🎯 Action (Δ³)",
         "• (α, β, γ) · 합 = 1\n"
         "• 연속 simplex\n"
         "• Dirichlet 분포에서 샘플링",
         ACCENT_RED),
        ("💰 Reward (scalar)",
         "R = 0.5·F1 + 0.3·EM\n"
         "      + 0.2·Faith\n"
         "      − 0.1·max(0, latency−5)",
         ACCENT_GREEN),
    ]
    left = Inches(0.8)
    w = Inches(3.9)
    for title, body, color in cards:
        _add_rounded_card(
            slide, left, Inches(2.7), w, Inches(2.5),
            title=title, body=body, title_color=color,
        )
        left += Inches(4.05)

    _add_text(
        slide, Inches(0.8), Inches(5.5),
        Inches(11.7), Inches(1.5),
        "에피소드 길이 = 1 — 가중치를 한 번 내면 retrieval·generation 이 일어나고\n"
        "즉시 보상이 결정되므로 다단계 transition 이 없음. 이 단순화가 학습을\n"
        "극히 안정적으로 만들고, 오프라인 보상 캐시 아이디어의 전제가 됨.",
        size=13, color=TEXT_DARK,
    )


def build_s09_network(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "신경망 구조 — Actor-Critic", "Network Architecture · 5,636 parameters")

    if (FIGS / "fig5_1_actor_critic.png").exists():
        slide.shapes.add_picture(
            str(FIGS / "fig5_1_actor_critic.png"),
            Inches(0.8), Inches(1.9),
            width=Inches(7.5),
        )

    _add_text(
        slide, Inches(8.6), Inches(2.0),
        Inches(4.3), Inches(0.5),
        "구성 요소",
        size=16, bold=True, color=HOSEO_BLUE,
    )
    _add_text(
        slide, Inches(8.6), Inches(2.6),
        Inches(4.3), Inches(4.0),
        "공유 백본:\n"
        "  18 → 64 → 64 (Tanh)\n\n"
        "Actor head:\n"
        "  64 → 3 (Dirichlet α_i)\n\n"
        "Critic head:\n"
        "  64 → 1 (V(s))\n\n"
        "총 파라미터 = 5,636\n"
        "(스마트폰 수준)",
        size=13, color=TEXT_DARK,
    )


def build_s10_ppo(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "학습 — PPO (Proximal Policy Optimization)", "Training with PPO")

    _add_text(
        slide, Inches(0.8), Inches(2.0),
        Inches(6), Inches(0.5),
        "학습 알고리즘",
        size=16, bold=True, color=HOSEO_BLUE,
    )
    _add_text(
        slide, Inches(0.8), Inches(2.6),
        Inches(6), Inches(4),
        "• OpenAI 2017, ChatGPT RLHF 에도 사용\n"
        "• clip ratio = 0.2 (정책 업데이트 폭 제한)\n"
        "• GAE λ = 0.95\n"
        "• Entropy coef = 0.01 (탐색 유지)\n"
        "• Total episodes = 10,000\n"
        "• 3 random seeds: 42, 123, 999",
        size=14, color=TEXT_DARK,
    )

    if (FIGS / "fig5_2_ppo_convergence.png").exists():
        slide.shapes.add_picture(
            str(FIGS / "fig5_2_ppo_convergence.png"),
            Inches(7.2), Inches(2.0),
            width=Inches(5.8),
        )

    _add_text(
        slide, Inches(0.8), Inches(6.5),
        Inches(11.7), Inches(0.7),
        "3-seed 표준편차 < 1% — 학습 자체의 재현성 확보.",
        size=13, color=HOSEO_BLUE, bold=True, align=PP_ALIGN.CENTER,
    )


def build_s11_cache(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "오프라인 보상 캐시 — 비용 병목 해소", "Offline Reward Cache")

    _add_text(
        slide, Inches(0.8), Inches(2.0),
        Inches(6), Inches(0.5),
        "문제: Naive PPO 는 매 에피소드 LLM 호출 필요",
        size=14, color=ACCENT_RED, bold=True,
    )
    _add_text(
        slide, Inches(0.8), Inches(2.5),
        Inches(6), Inches(1.5),
        "• 10,000 ep × 3 seed × 32 rollout\n"
        "  = 약 1M LLM 호출\n"
        "• gpt-4o-mini 단가: ≈ $288\n"
        "• 개인 연구자가 반복 실험 불가",
        size=13, color=TEXT_DARK,
    )
    _add_text(
        slide, Inches(0.8), Inches(4.0),
        Inches(6), Inches(0.5),
        "해결: 보상을 미리 전부 계산",
        size=14, color=ACCENT_GREEN, bold=True,
    )
    _add_text(
        slide, Inches(0.8), Inches(4.5),
        Inches(6), Inches(2.3),
        "① 5,000 질의 × 66 이산 조합 = 330,000 엔트리\n"
        "② (α, β, γ) → (F1, EM, Faith, latency) 를 SQLite 에 저장\n"
        "③ 이후 PPO 학습은 조회만 — LLM 호출 0 회\n\n"
        "캐시 구축 1회: $33 · 14시간\n"
        "이후 3-seed 학습: 5시간 · $0\n"
        "→ 총 비용 −85% 절감",
        size=13, color=TEXT_DARK,
    )

    # Cost comparison card
    _add_rounded_card(
        slide, Inches(7.2), Inches(2.3), Inches(5.8), Inches(3.8),
        title="비용 구조 비교",
        body="              Naive PPO         Cache + PPO\n"
             "─────────────────────────────────────\n"
             "LLM 호출       ~1,000,000          330,000 (1회)\n"
             "학습 호출      10,000,000+                  0\n"
             "총 LLM 비용        $288                    $33\n"
             "wall-clock         50 h+                  19 h\n"
             "재학습 비용      매번 $288             매번 $0\n",
        fill_color=BG_SOFT,
    )

    _add_text(
        slide, Inches(7.2), Inches(6.3),
        Inches(5.8), Inches(0.8),
        "→ 이 엔지니어링 선택이 개인 연구자 수준의\n"
        "   반복 가능성을 만들어 주었다.",
        size=13, color=HOSEO_BLUE, bold=True,
    )


def build_s12_corrigendum(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "재현성 복구 — Stage-wise CORRIGENDUM",
                        "Reproducibility Recovery (S0 → S3)")

    _add_text(
        slide, Inches(0.8), Inches(2.0),
        Inches(11.7), Inches(0.5),
        "선행 JKSCI 2025 의 F1 0.86 이 현 저장소로는 재현 불가",
        size=16, bold=True, color=ACCENT_RED,
    )

    stages = [
        ("S0", "JKSCI 원 주장", "F1 ~0.86", "재현 불가 (CORRIGENDUM §3)", TEXT_MUTED),
        ("S1", "`normalize_korean` 구두점 버그 수정",
         "F1 0.137", "토큰 분리 오류 → +90% 상승", HOSEO_BLUE),
        ("S2", "`PROMPT_TEMPLATE_LIST` 로 gold 형식 정렬",
         "F1 0.529", "프롬프트 리스트 형식 → +286%", HOSEO_BLUE_LIGHT),
        ("S3", "`faithfulness()` 2-branch 도입",
         "F1 0.529", "리스트형 엄격화 (측정 정직성)", ACCENT_GREEN),
    ]
    top = Inches(2.7)
    for tag, desc, f1, note, color in stages:
        _add_rect(slide, Inches(0.8), top, Inches(0.6), Inches(0.7),
                  fill_color=color)
        _add_text(slide, Inches(0.8), top + Inches(0.1),
                  Inches(0.6), Inches(0.5),
                  tag, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_text(slide, Inches(1.5), top,
                  Inches(6.5), Inches(0.4),
                  desc, size=13, bold=True, color=TEXT_DARK)
        _add_text(slide, Inches(1.5), top + Inches(0.35),
                  Inches(6.5), Inches(0.4),
                  note, size=11, color=TEXT_MUTED)
        _add_text(slide, Inches(8.3), top + Inches(0.15),
                  Inches(4.3), Inches(0.5),
                  f1, size=16, bold=True, color=color, align=PP_ALIGN.RIGHT)
        top += Inches(0.95)

    _add_text(
        slide, Inches(0.8), Inches(6.5),
        Inches(11.7), Inches(0.7),
        "본 논문은 S1+S2+S3 적용된 corrected baseline 위에서 L-DWA 의 순수 학습 기여를 측정.",
        size=13, color=HOSEO_BLUE, bold=True, align=PP_ALIGN.CENTER,
    )


def build_s13_experiment(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "실험 설계", "Experimental Setup")

    items = [
        ("📚 벤치마크",
         "합성 대학 코퍼스 (자체 구축)\n"
         "60 학과 / 577 교수 / 1,505 과목\n"
         "/ 400 프로젝트 / 2,542 문서\n\n"
         "Gold QA: 5,000 쌍\n"
         "(simple 2,000 / multi_hop 1,750\n"
         "  / conditional 1,250)"),
        ("🔧 LLM · Embedding",
         "gpt-4o-mini (2024-07-18)\n"
         "temperature = 0.0\n"
         "max_tokens = 500\n\n"
         "text-embedding-3-small\n"
         "(1536-dim)\n"
         "FAISS IndexFlatIP, top-k=3"),
        ("📏 평가 지표",
         "F1_strict (엄격 토큰 F1)\n"
         "F1_substring (부분 일치 F1)\n"
         "F1_char (문자 3-gram F1)\n"
         "EM (Exact Match)\n"
         "Faithfulness (2-branch)\n"
         "Latency (seconds)"),
    ]
    left = Inches(0.8)
    w = Inches(3.9)
    for title, body in items:
        _add_rounded_card(
            slide, left, Inches(2.0), w, Inches(4.5),
            title=title, body=body, fill_color=BG_SOFT,
        )
        left += Inches(4.05)

    _add_text(
        slide, Inches(0.8), Inches(6.7),
        Inches(11.7), Inches(0.5),
        "3 random seeds (42, 123, 999) 평균 보고 · post-CORRIGENDUM 평가기",
        size=12, color=TEXT_MUTED, align=PP_ALIGN.CENTER,
    )


def build_s14_results_table(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "전체 결과 — 정책별 비교 (5,000 QA)", "Overall Comparison")

    headers = ["정책", "F1_strict", "F1_substring", "F1_char", "EM", "Faith", "Latency"]
    rows = [
        ["Vector-only", "0.334", "0.334", "0.317", "0.250", "0.470", "0.68s", TEXT_MUTED],
        ["R-DWA", "0.529", "0.482", "0.469", "0.387", "0.544", "0.71s", HOSEO_BLUE],
        ["L-DWA (3-seed)", "0.562", "0.507", "0.494", "0.388", "0.580", "0.76s", ACCENT_RED],
        ["Oracle", "0.554", "0.504", "0.487", "0.388", "0.570", "0.75s", ACCENT_GREEN],
    ]

    table_left = Inches(0.8)
    table_top = Inches(2.0)
    col_w = [Inches(2.3)] + [Inches(1.5)] * 6
    row_h = Inches(0.55)

    # Header row
    x = table_left
    for i, h in enumerate(headers):
        _add_rect(slide, x, table_top, col_w[i], row_h, fill_color=HOSEO_BLUE)
        _add_text(
            slide, x, table_top + Inches(0.12),
            col_w[i], Inches(0.4),
            h, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        )
        x += col_w[i]

    y = table_top + row_h
    for row in rows:
        x = table_left
        color = row[-1]
        cells = row[:-1]
        for i, val in enumerate(cells):
            _add_rect(slide, x, y, col_w[i], row_h,
                      fill_color=BG_SOFT, line_color=HOSEO_BLUE_LIGHT)
            _add_text(
                slide, x, y + Inches(0.14),
                col_w[i], Inches(0.4),
                val,
                size=13, bold=(i == 0),
                color=color if i == 0 else TEXT_DARK,
                align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER,
            )
            x += col_w[i]
        y += row_h

    _add_text(
        slide, Inches(0.8), Inches(5.4),
        Inches(11.7), Inches(1.6),
        "핵심: L-DWA 3-seed 평균이 R-DWA 대비 F1_strict +6.2%,\n"
        "       Faithfulness +6.6%, 그리고 Oracle 0.554 를 네 F1 축 모두에서 초과.\n"
        "       3-seed 표준편차는 ±0.007 수준으로 극히 안정적.",
        size=14, color=HOSEO_BLUE, bold=True,
    )


def build_s15_per_type(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "유형별 F1 — 어디서 개선 폭이 큰가", "Per-Type Breakdown")

    if (FIGS / "fig6_2_per_type.png").exists():
        slide.shapes.add_picture(
            str(FIGS / "fig6_2_per_type.png"),
            Inches(0.8), Inches(1.9),
            width=Inches(8.0),
        )

    _add_text(
        slide, Inches(9.0), Inches(2.0),
        Inches(4.0), Inches(0.5),
        "유형별 F1_strict",
        size=16, bold=True, color=HOSEO_BLUE,
    )
    _add_text(
        slide, Inches(9.0), Inches(2.6),
        Inches(4.0), Inches(4),
        "Simple (n=2000):\n"
        "  R-DWA   0.874\n"
        "  L-DWA  0.906\n"
        "  Oracle  0.901\n\n"
        "Multi-hop (n=1750):\n"
        "  R-DWA   0.354\n"
        "  L-DWA  0.365\n"
        "  Oracle  0.380\n\n"
        "Conditional (n=1250): ★\n"
        "  R-DWA   0.223\n"
        "  L-DWA  0.304  (+36.7%)\n"
        "  Oracle  0.290",
        size=12, color=TEXT_DARK, font=EN_FONT,
    )


def build_s16_oracle_exceedance(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "핵심 결과 ① — 이산 Oracle 상한 초과",
                        "Key Finding 1 · Continuous Policy Exceeds Discrete Oracle")

    _add_text(
        slide, Inches(0.8), Inches(2.0),
        Inches(11.7), Inches(1.2),
        "66-점 이산 격자 argmax Oracle (F1_strict 0.554) 을\n"
        "학습된 연속 Dirichlet 평균이 네 개 F1 축 모두에서 미세하게 초과",
        size=18, bold=True, color=ACCENT_RED,
    )

    items = [
        ("F1_strict",       "0.529", "0.554", "0.562"),
        ("F1_substring",    "0.482", "0.504", "0.507"),
        ("F1_char",         "0.469", "0.487", "0.494"),
        ("Faithfulness",    "0.544", "0.570", "0.580"),
    ]
    top = Inches(3.4)
    # header
    x = Inches(0.8)
    col_w = [Inches(3.0), Inches(2.8), Inches(2.8), Inches(2.8)]
    hdr = ["지표", "R-DWA", "Oracle", "L-DWA"]
    for i, h in enumerate(hdr):
        _add_rect(slide, x, top, col_w[i], Inches(0.5),
                  fill_color=HOSEO_BLUE)
        _add_text(slide, x, top + Inches(0.1),
                  col_w[i], Inches(0.4),
                  h, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        x += col_w[i]
    y = top + Inches(0.5)
    for metric, r, o, ld in items:
        x = Inches(0.8)
        for i, val in enumerate([metric, r, o, ld]):
            c = BG_SOFT if i > 0 else WHITE
            _add_rect(slide, x, y, col_w[i], Inches(0.45),
                      fill_color=c, line_color=HOSEO_BLUE_LIGHT)
            color = ACCENT_RED if i == 3 else TEXT_DARK
            bold = (i == 3)
            _add_text(slide, x, y + Inches(0.08),
                      col_w[i], Inches(0.4),
                      val, size=14, bold=bold, color=color,
                      align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER)
            x += col_w[i]
        y += Inches(0.45)

    _add_text(
        slide, Inches(0.8), Inches(6.3),
        Inches(11.7), Inches(0.9),
        "해석 — Oracle 은 66개 이산 점 중 argmax. L-DWA 의 Dirichlet 평균은\n"
        "연속 공간 → 격자 바깥 지점을 활용 가능. 양적 개선이 아닌 질적 의미.",
        size=13, color=HOSEO_BLUE,
    )


def build_s17_conditional(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "핵심 결과 ② — Conditional +36.7%",
                        "Key Finding 2 · Largest Per-Type Gain in Conditional")

    _add_text(
        slide, Inches(0.8), Inches(2.0),
        Inches(11.7), Inches(1.2),
        "조건부 질의 (n = 1,250) 에서 가장 큰 상대 개선.\n"
        "L-DWA 가 Oracle 마저 근소하게 초과.",
        size=18, bold=True, color=ACCENT_RED,
    )

    # 3 metric cards
    cards = [
        ("R-DWA\n(기존 규칙 표)", "0.223", TEXT_MUTED, "고정 (0.2, 0.2, 0.6)"),
        ("Oracle\n(이산 상한)", "0.290", ACCENT_GREEN, "66-점 argmax"),
        ("L-DWA\n(본 연구)", "0.304", ACCENT_RED, "질의별 동적 가중치"),
    ]
    left = Inches(0.8)
    w = Inches(3.9)
    for label, val, color, sub in cards:
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     left, Inches(3.4), w, Inches(2.5))
        shp.fill.solid()
        shp.fill.fore_color.rgb = WHITE
        shp.line.color.rgb = color
        shp.line.width = Pt(2.5)
        _add_text(slide, left, Inches(3.55),
                  w, Inches(0.8),
                  label, size=14, bold=True,
                  color=color, align=PP_ALIGN.CENTER)
        _add_text(slide, left, Inches(4.3),
                  w, Inches(1.0),
                  val, size=42, bold=True,
                  color=color, align=PP_ALIGN.CENTER)
        _add_text(slide, left, Inches(5.3),
                  w, Inches(0.5),
                  sub, size=11, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
        left += Inches(4.05)

    _add_text(
        slide, Inches(0.8), Inches(6.2),
        Inches(11.7), Inches(1.0),
        "해석 — 같은 conditional 이라도 수치 제약 · 열거형 · 배타 조건이 각기 다른데,\n"
        "R-DWA 는 이들을 하나로 묶어 처리. 학습이 내부 차이를 포착.",
        size=13, color=HOSEO_BLUE,
    )


def build_s18_weight_distribution(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "Δ³ 위 정책별 평균 가중치 분포",
                        "Mean Weight Positions on the 3-Simplex")

    if (FIGS / "fig6_3_weight_distribution.png").exists():
        slide.shapes.add_picture(
            str(FIGS / "fig6_3_weight_distribution.png"),
            Inches(0.8), Inches(1.9),
            width=Inches(7.5),
        )

    _add_text(
        slide, Inches(8.6), Inches(2.0),
        Inches(4.3), Inches(0.5),
        "읽는 법",
        size=16, bold=True, color=HOSEO_BLUE,
    )
    _add_text(
        slide, Inches(8.6), Inches(2.6),
        Inches(4.3), Inches(4.0),
        "R-DWA 평균 (파랑):\n"
        "  (0.25, 0.45, 0.30)\n"
        "  Graph 쪽 치우침\n\n"
        "L-DWA 평균 (빨강):\n"
        "  (0.25, 0.30, 0.45)\n"
        "  Ontology 쪽 이동\n\n"
        "Oracle 평균 (초록):\n"
        "  (0.11, 0.16, 0.73)\n"
        "  Ontology 꼭짓점 근처",
        size=12, color=TEXT_DARK, font=EN_FONT,
    )

    _add_text(
        slide, Inches(0.8), Inches(6.5),
        Inches(11.7), Inches(0.7),
        "정답 분포의 진짜 지형은 Ontology 지배적. R-DWA 는 이를 반영하지 못함.",
        size=13, color=HOSEO_BLUE, bold=True, align=PP_ALIGN.CENTER,
    )


def build_s19_ppo_curves(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "PPO 3-seed 학습 곡선",
                        "PPO Convergence across 3 Seeds")

    if (FIGS / "fig5_2_ppo_convergence.png").exists():
        slide.shapes.add_picture(
            str(FIGS / "fig5_2_ppo_convergence.png"),
            Inches(1.3), Inches(1.9),
            width=Inches(10.7),
        )

    _add_text(
        slide, Inches(0.8), Inches(6.2),
        Inches(11.7), Inches(1.0),
        "세 시드 모두 mean_reward ≈ 0.215 ± 0.002 로 수렴.\n"
        "1% 미만의 표준편차 → 학습 자체의 재현 가능성 확보.",
        size=14, color=HOSEO_BLUE, bold=True, align=PP_ALIGN.CENTER,
    )


def build_s20_stagewise(prs: Presentation) -> None:
    """Re-use S12 CORRIGENDUM style but for Results section - maybe skip."""
    # Alternative: do cross-domain here instead
    pass


def build_s21_cross_domain(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "교차 도메인 — naive transfer 의 한계와 원인 분해",
                        "Cross-domain · Failure Decomposed into 3 Axes")

    if (FIGS / "fig6_5_cross_domain_radar.png").exists():
        slide.shapes.add_picture(
            str(FIGS / "fig6_5_cross_domain_radar.png"),
            Inches(0.8), Inches(1.9),
            width=Inches(6.5),
        )

    _add_text(
        slide, Inches(7.5), Inches(2.0),
        Inches(5.3), Inches(0.5),
        "실패의 세 층위",
        size=16, bold=True, color=ACCENT_RED,
    )
    _add_text(
        slide, Inches(7.5), Inches(2.6),
        Inches(5.3), Inches(4.3),
        "① 언어 장벽\n"
        "    한국어 정규식 intent → 영어 질의 density ≈ 0\n"
        "    → 영어 intent 추가 시 R-DWA +38%, L-DWA +49%\n\n"
        "② 도메인 어휘\n"
        "    학습된 엔티티 vs HotpotQA 엔티티 상이\n"
        "    → 영어 합성 벤치마크 R-DWA 0.663\n"
        "       (한국어 0.529 보다 높음)\n\n"
        "③ Graph·Ontology 부재\n"
        "    HotpotQA 는 그래프 구조 없음\n"
        "    → 재학습 해도 곡선 평탄\n"
        "    = Triple-Hybrid 의 본질적 경계",
        size=12, color=TEXT_DARK,
    )


def build_s23_contributions(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "이 논문의 기여 네 가지", "Four Contributions")

    cards = [
        ("방법론", "Triple-Hybrid RAG 의 가중치 결정을 1-step MDP 로 공식화하고 PPO 로 학습한 최초 사례",
         HOSEO_BLUE),
        ("엔지니어링", "오프라인 보상 캐시 (330K, 16.8 MB) 로 학습 비용 $288 → $33 (−85%).\n"
                     "RL-based RAG 튜닝을 개인 연구자 수준으로 접근 가능하게",
         HOSEO_BLUE_LIGHT),
        ("재현성", "선행 JKSCI 2025 의 평가 코드 결함 3건 발견 · 수정 · CORRIGENDUM 공개 정정\n"
                "(구두점 처리 · 프롬프트 형식 · Faithfulness 2-branch)",
         ACCENT_GREEN),
        ("실증", "L-DWA 가 R-DWA 대비 F1_strict +6.2%, 이산 Oracle 상한 0.554 를 0.562 로 초과.\n"
               "Conditional 유형에서 +36.7% 의 가장 큰 상대 개선.",
         ACCENT_RED),
    ]
    top = Inches(2.0)
    for label, body, color in cards:
        _add_rect(slide, Inches(0.8), top, Inches(0.3), Inches(1.0),
                  fill_color=color)
        _add_text(slide, Inches(1.2), top + Inches(0.05),
                  Inches(2.2), Inches(0.5),
                  label, size=16, bold=True, color=color)
        _add_text(slide, Inches(3.5), top + Inches(0.05),
                  Inches(9.2), Inches(1.0),
                  body, size=13, color=TEXT_DARK)
        top += Inches(1.15)


def build_s24_limitations(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "한계 및 향후 연구 방향",
                        "Limitations & Future Work")

    _add_text(
        slide, Inches(0.8), Inches(2.0),
        Inches(6), Inches(0.5),
        "솔직하게 보고하는 한계",
        size=16, bold=True, color=ACCENT_RED,
    )
    _add_text(
        slide, Inches(0.8), Inches(2.6),
        Inches(6), Inches(4),
        "• 한국어 단일 도메인 end-to-end 학습\n\n"
        "• State 품질이 L-DWA 상한 결정\n"
        "    (영어 state 미구분 → R-DWA 수준)\n\n"
        "• EM 구조적 문제 — list prompt 로 0.39 회복\n\n"
        "• BERT intent classifier 미학습\n\n"
        "• 교차 도메인 naive transfer −30~35%",
        size=13, color=TEXT_DARK,
    )

    _add_text(
        slide, Inches(7.2), Inches(2.0),
        Inches(5.8), Inches(0.5),
        "향후 연구 방향",
        size=16, bold=True, color=HOSEO_BLUE,
    )
    _add_text(
        slide, Inches(7.2), Inches(2.6),
        Inches(5.8), Inches(4),
        "• 다국어 state feature 설계\n\n"
        "• 도메인별 fine-tuning 레시피 체계화\n\n"
        "• Reward 재설계 (EM 가중치 재분배)\n\n"
        "• Joint multi-domain training\n\n"
        "• 더 큰 Actor-Critic 의 스케일링 효과\n\n"
        "• 선행 저장소와 평가기 통일",
        size=13, color=TEXT_DARK,
    )


def build_s25_lessons(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "연구 과정에서 얻은 두 교훈",
                        "Two Lessons Learned")

    _add_rounded_card(
        slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.2),
        title="① 평가 코드를 다시 들여다본 경험",
        body="마무리 단계에서 F1 수치가 기대보다 낮다는 인상을 받아 "
             "normalize_korean 을 역추적한 끝에, 쉼표와 마침표가 토큰화 "
             "이전에 제거되지 않는 작은 버그를 발견. 이를 바로잡고 나서야 "
             "L-DWA 의 Oracle 대비 위치가 제대로 드러남.  \n"
             "→ 실험 결과가 직관과 어긋날 때 평가 코드부터 다시 읽어 보는 습관의 중요성.",
        title_color=HOSEO_BLUE,
    )

    _add_rounded_card(
        slide, Inches(0.8), Inches(4.4), Inches(11.7), Inches(2.2),
        title="② 기대와 다르게 나온 결과를 어떻게 다룰지",
        body="교차 도메인 실험이 30~35% 저하로 나오자 한 문단 \"한계\" 로 "
             "넘기고 싶은 유혹이 있었음. 그러나 원인을 언어 · 도메인 · 아키텍처로 "
             "나누어 각각 별도 실험으로 확인해 보는 편이 오히려 더 명확한 결론을 냄.  \n"
             "→ 부정적 결과를 빨리 결론짓지 않는 쪽이 장기적으로 나은 서사를 만든다.",
        title_color=ACCENT_GREEN,
    )


def build_s26_closing(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    # Full screen blue panel
    _add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=HOSEO_BLUE)

    _add_text(
        slide, Inches(0), Inches(2.0),
        SLIDE_W, Inches(0.8),
        "감사합니다",
        size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide, Inches(0), Inches(3.0),
        SLIDE_W, Inches(0.6),
        "Thank you · 질의응답 (Q&A)",
        size=20, color=WHITE, align=PP_ALIGN.CENTER, font=EN_FONT,
    )

    _add_text(
        slide, Inches(0), Inches(4.5),
        SLIDE_W, Inches(0.5),
        "공개 자원",
        size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide, Inches(0), Inches(5.0),
        SLIDE_W, Inches(1.8),
        "GitHub: github.com/sdw1621/triple-rag-phd\n"
        "Dashboard: triple-rag-phd.streamlit.app\n"
        "CORRIGENDUM: github.com/sdw1621/hybrid-rag-comparsion/blob/main/CORRIGENDUM.md\n\n"
        "신동욱 · 호서대학교 대학원 융합공학과 · 2026",
        size=14, color=WHITE, align=PP_ALIGN.CENTER, font=EN_FONT,
    )


# --------------------------------------------------------------------------
# Appendix
# --------------------------------------------------------------------------

def build_appendix_hp(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "Appendix A — PPO Hyperparameters",
                        "Table 5-4 reproduction")

    items = [
        ("learning_rate", "3.0e-4"),
        ("gae_lambda", "0.95"),
        ("clip_ratio", "0.2"),
        ("value_coef", "0.5"),
        ("entropy_coef", "0.01"),
        ("max_grad_norm", "0.5"),
        ("total_episodes", "10,000"),
        ("rollout_per_episode", "32"),
        ("update_epochs", "4"),
        ("minibatch_size", "8"),
        ("gamma", "0.99"),
        ("seeds", "42, 123, 999"),
    ]
    left = Inches(0.8)
    top = Inches(2.0)
    col_w = Inches(4.0)
    row_h = Inches(0.45)
    for i, (k, v) in enumerate(items):
        col_idx = i // 6
        row_idx = i % 6
        x = left + col_idx * Inches(6.3)
        y = top + row_idx * row_h
        _add_rect(slide, x, y, col_w, row_h,
                  fill_color=BG_SOFT, line_color=HOSEO_BLUE_LIGHT)
        _add_rect(slide, x + col_w, y, Inches(2.2), row_h,
                  fill_color=WHITE, line_color=HOSEO_BLUE_LIGHT)
        _add_text(slide, x + Inches(0.15), y + Inches(0.1),
                  col_w, Inches(0.4),
                  k, size=12, bold=True, color=TEXT_DARK, font=EN_FONT)
        _add_text(slide, x + col_w, y + Inches(0.1),
                  Inches(2.0), Inches(0.4),
                  v, size=12, color=HOSEO_BLUE, font=EN_FONT, align=PP_ALIGN.CENTER)


def build_appendix_cases(prs: Presentation) -> None:
    slide = _blank_slide(prs)
    _set_bg(slide, WHITE)
    _add_section_header(slide, "Appendix B — Case Studies CS-1 ~ CS-7 요약",
                        "Qualitative case summaries")

    cases = [
        ("CS-1", "Simple 정책 수렴", "세 정책 모두 F1=1.0 (최재원 교수 소속 학과)"),
        ("CS-2", "Multi-hop 부분 회수", "R-DWA/L-DWA 동일 1/2 엔티티, Oracle 2/2"),
        ("CS-3", "Conditional 개선과 한계", "개별 실패지만 유형 집계 +36.7% (L-DWA > Oracle)"),
        ("CS-4", "Faithfulness 2-branch", "Sentence 1.0 (과대) vs List 0.5 (환각 노출)"),
        ("CS-5", "Retrieval 한계", "Oracle 조차 12개 중 2개만 — 컨텍스트 상한"),
        ("CS-6", "교차 도메인 퇴행", "영어 질의 → density 0 → Vector-only 이하"),
        ("CS-7", "List 프롬프트 효과", "F1_strict 0.137 → 0.529 (×3.86), 평균 latency −50%"),
    ]
    top = Inches(2.0)
    for tag, title, body in cases:
        _add_rect(slide, Inches(0.8), top, Inches(1.0), Inches(0.55),
                  fill_color=HOSEO_BLUE)
        _add_text(slide, Inches(0.8), top + Inches(0.1),
                  Inches(1.0), Inches(0.4),
                  tag, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        _add_text(slide, Inches(2.0), top + Inches(0.05),
                  Inches(3.5), Inches(0.45),
                  title, size=13, bold=True, color=TEXT_DARK)
        _add_text(slide, Inches(5.5), top + Inches(0.1),
                  Inches(7.2), Inches(0.5),
                  body, size=12, color=TEXT_MUTED)
        top += Inches(0.65)


# --------------------------------------------------------------------------
# main assembler
# --------------------------------------------------------------------------

SLIDES = [
    ("Part A — Opening", [
        ("title", build_s01_title),
        ("agenda", build_s02_agenda),
        ("summary", build_s03_summary),
    ]),
    ("Part B — 문제와 방법", [
        ("bg_rag", build_s04_rag_bg),
        ("triple_hybrid", build_s05_triple_hybrid),
        ("weight_problem", build_s06_weight_problem),
        ("rdwa", build_s07_rdwa),
        ("ldwa_mdp", build_s08_ldwa_mdp),
        ("network", build_s09_network),
        ("ppo", build_s10_ppo),
        ("cache", build_s11_cache),
        ("corrigendum", build_s12_corrigendum),
    ]),
    ("Part C — 결과", [
        ("experiment", build_s13_experiment),
        ("results_table", build_s14_results_table),
        ("per_type", build_s15_per_type),
        ("oracle_exceedance", build_s16_oracle_exceedance),
        ("conditional", build_s17_conditional),
        ("weight_dist", build_s18_weight_distribution),
        ("ppo_curves", build_s19_ppo_curves),
        ("cross_domain", build_s21_cross_domain),
    ]),
    ("Part D — 결론", [
        ("contributions", build_s23_contributions),
        ("limitations", build_s24_limitations),
        ("lessons", build_s25_lessons),
        ("closing", build_s26_closing),
    ]),
    ("Appendix", [
        ("hp", build_appendix_hp),
        ("cases", build_appendix_cases),
    ]),
]


def build_all(output: Path) -> None:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # flatten
    flat: list[tuple[str, str, callable]] = []
    for section, entries in SLIDES:
        for tag, fn in entries:
            flat.append((section, tag, fn))
    total = len(flat)

    for idx, (section, tag, fn) in enumerate(flat, start=1):
        fn(prs)
        slide = prs.slides[idx - 1]
        if tag != "title" and tag != "closing":
            _add_footer(slide, idx, total, section=section)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f"Wrote {output}  ({total} slides)")


def main() -> int:
    p = argparse.ArgumentParser()
    p.add_argument(
        "--output",
        default=str(ROOT / "thesis_current" / "박사논문_방어발표_v1.pptx"),
    )
    args = p.parse_args()
    build_all(Path(args.output))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
